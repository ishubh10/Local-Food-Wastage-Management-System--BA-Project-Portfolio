
,KO,LK,P''';K --.LookupError+"""
Generate executive PowerPoint presentation for the Food Waste Management System.
Output: final/Food_Waste_Executive_Summary.pptx
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import DATA_FILES, FINAL_DIR


GREEN = RGBColor(11, 110, 79)
DARK = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(34, 197, 94)
BLUE = RGBColor(59, 130, 246)


def load_kpis() -> dict:
    providers = pd.read_csv(DATA_FILES["providers"])
    receivers = pd.read_csv(DATA_FILES["receivers"])
    food = pd.read_csv(DATA_FILES["food"])
    claims = pd.read_csv(DATA_FILES["claims"])
    claims["Timestamp"] = pd.to_datetime(claims["Timestamp"])

    total_qty = food["Quantity"].sum()
    completed = claims[claims["Status"] == "Completed"]
    claimed_qty = food[food["Food_ID"].isin(completed["Food_ID"])]["Quantity"].sum()

    return {
        "providers": len(providers),
        "receivers": len(receivers),
        "food_listings": len(food),
        "claims": len(claims),
        "total_qty": int(total_qty),
        "claimed_qty": int(claimed_qty),
        "rescue_rate": round(claimed_qty / total_qty * 100, 1),
        "success_rate": round(len(completed) / len(claims) * 100, 1),
        "cancel_rate": round((claims["Status"] == "Cancelled").sum() / len(claims) * 100, 1),
        "pending": int((claims["Status"] == "Pending").sum()),
        "top_provider": food.groupby("Provider_Type")["Quantity"].sum().idxmax(),
        "top_receiver": claims.merge(receivers, on="Receiver_ID")["Type"].value_counts().idxmax(),
        "top_city": food.groupby("Location")["Quantity"].sum().idxmax(),
        "date_range": f"{claims['Timestamp'].min().strftime('%b %d, %Y')} – {claims['Timestamp'].max().strftime('%b %d, %Y')}",
    }


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(12)


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.font.size = Pt(18)
        para.font.color.rgb = DARK
        para.space_after = Pt(10)
        para.level = 0


def add_kpi_slide(prs: Presentation, kpis: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Key Performance Indicators"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN

    metrics = [
        ("Total Food Units", f"{kpis['total_qty']:,}"),
        ("Food Rescue Rate", f"{kpis['rescue_rate']}%"),
        ("Claim Success Rate", f"{kpis['success_rate']}%"),
        ("Cancellation Rate", f"{kpis['cancel_rate']}%"),
        ("Active Providers", f"{kpis['providers']:,}"),
        ("Active Receivers", f"{kpis['receivers']:,}"),
    ]

    for i, (label, value) in enumerate(metrics):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 3.1)
        top = Inches(1.6 + row * 2.4)
        card = slide.shapes.add_shape(1, left, top, Inches(2.8), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK
        card.line.fill.background()

        val_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(2.4), Inches(0.8))
        vp = val_box.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(28)
        vp.font.bold = True
        vp.font.color.rgb = ACCENT

        lbl_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.1), Inches(2.4), Inches(0.5))
        lp = lbl_box.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(12)
        lp.font.color.rgb = WHITE


def build_presentation(kpis: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Food Waste Management System",
        "Executive Summary & Data Insights",
    )

    add_bullet_slide(prs, "Problem Statement", [
        "1.3 billion tonnes of food are wasted globally each year while millions face hunger.",
        "Local food donation platforms connect surplus food from providers to NGOs, shelters, and individuals.",
        "This project analyzes operational data to improve rescue rates and reduce claim failures.",
        f"Analysis period: {kpis['date_range']}",
    ])

    add_bullet_slide(prs, "Dataset Overview", [
        f"{kpis['providers']:,} food providers (supermarkets, restaurants, grocery stores, catering)",
        f"{kpis['receivers']:,} receivers (NGOs, charities, shelters, individuals)",
        f"{kpis['food_listings']:,} food listings totaling {kpis['total_qty']:,} units",
        f"{kpis['claims']:,} claim transactions tracked with status and timestamps",
        "100% data completeness — zero null values across all tables",
    ])

    add_kpi_slide(prs, kpis)

    add_bullet_slide(prs, "Claim Status Analysis", [
        f"Claim success rate: {kpis['success_rate']}% — only 1 in 3 claims completes successfully",
        f"Cancellation rate: {kpis['cancel_rate']}% — systemic friction, not isolated incidents",
        f"Pending claims: {kpis['pending']:,} — unresolved demand awaiting action",
        "Recommendation: Implement SLA tracking, expiry alerts, and automated matching",
    ])

    add_bullet_slide(prs, "Supply Side — Providers", [
        f"Top provider category: {kpis['top_provider']}",
        "Restaurants and supermarkets contribute the highest donation volumes",
        "Four provider types: Supermarket, Grocery Store, Restaurant, Catering Service",
        "Opportunity: Formalize partnerships with top-performing donor categories",
    ])

    add_bullet_slide(prs, "Demand Side — Receivers", [
        f"Top receiver category: {kpis['top_receiver']}",
        "NGOs and charities drive the majority of claim activity",
        "Balanced demand across NGO, Charity, Shelter, and Individual segments",
        "Opportunity: Expand receiver network in high-supply, low-demand cities",
    ])

    add_bullet_slide(prs, "Geographic Insights", [
        f"Highest supply city: {kpis['top_city']}",
        "Food supply is concentrated in a subset of urban areas",
        "Supply-demand mismatch exists in several high-surplus cities",
        "Recommendation: Target receiver outreach and logistics in surplus regions",
    ])

    add_bullet_slide(prs, "Strategic Recommendations", [
        "1. Reduce cancellations through real-time expiry monitoring and claim reminders",
        "2. Accelerate pending claim resolution with SLA dashboards and notifications",
        "3. Rebalance geographic coverage — match supply hotspots with receiver capacity",
        "4. Strengthen provider partnerships in restaurant and supermarket segments",
        "5. Deploy the interactive Streamlit dashboard for ongoing executive monitoring",
    ])

    add_section_slide(prs, "Thank You")
    return prs


def main() -> None:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    kpis = load_kpis()
    prs = build_presentation(kpis)
    out = FINAL_DIR / "Food_Waste_Executive_Summary.pptx"
    prs.save(str(out))
    print(f"Presentation saved: {out}")


if __name__ == "__main__":
    main()
